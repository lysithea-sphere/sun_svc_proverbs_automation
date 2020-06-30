from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.dml import MSO_COLOR_TYPE
from pptx.enum.text import MSO_AUTO_SIZE
from pptx.enum.dml import MSO_THEME_COLOR
from pptx.dml.color import ColorFormat, RGBColor
import os
from string_breaker import eng_str_break, chi_str_break, list_string_remover

prs = Presentation(r"C:\Users\Wei Jie\PycharmProjects\H_World\venv\Proverb_proj\Sun_Svc_Proverbs_automation\Dependencies\alpha.pptx")  # Path or Filename ***IMPORTANT**

eng_filename = ('eng_proverbs.txt')     #specify directory or cd into it, where default is reading into current root directory
with open(eng_filename, encoding='utf8') as f:
    eng_text = f.read()
eng_list = eng_text.splitlines()

chi_filename = ('chi_proverbs.txt')     #specify directory or cd into it, where default is reading into current root directory
with open(chi_filename, encoding='utf8') as f:
    chi_text = f.read()
chi_text = chi_text.replace('　', '')        #this is needed due to an extra character found that is causing line break to occur. character : '　'
chi_list = chi_text.splitlines()

full_list = []
count = 0
while count < 12:
    full_list.append(eng_list[count])
    full_list.append(chi_list[count])
    count += 2
#full_list is now eng_proverb, chi_proverb alternate, but excludes 'together'.
is_english = False      #preset to False
is_presider = False     #preset to False

eng_limit = 108  # limit per slide for eng; A single line ranges from 38-46 characters, thus, being prudent and accounting for worst scenario, 36 per line is used, where 36 x 3 = 108
chi_limit = 63  # 21 characters per line; limit per slide (chinese auto monospace?)
#full_list is fully correct till here

for proverb in full_list:
    if 'Presider' in proverb or '司會' in proverb or '司会' in proverb:
        is_presider = True
    else:
        is_presider = False

    if 'Presider' in proverb or 'Congregation' in proverb:
        is_english = True
    else:
        is_english = False
    #is_presider and is_english is now determined by the conditional statements above

    #if is_presider and is_english:
    pro = proverb.replace('Presider', '', 1)      #to remove the 'Presider' demarkation at the front
    #elif is_presider and not is_english :   #to remove the '司会' demarkation at the front
    pro = pro.replace('司会', '', 1)
    pro = pro.replace('司會', '', 1)
    #elif not is_presider and is_english:
    pro = pro.replace('Congregation', '', 1)  #to remove the 'Congregation' demarkation at the front
    #else:
    pro = pro.replace('会众', '', 1)          #to remove '会众' demarkation at the front
    pro = pro.replace('會眾', '', 1)

    pro = pro.replace('-', '', 1)                 #to remove the hyphen('-') in front
    pro = pro.replace('－', '', 1)                #to remove the chinese hyphen('－') in front
    pro = pro.lstrip()                            #strips the extra spaces in front
    pro = pro.rstrip()                            #strips the extra spaces at the end so that character count is accurate

    if is_english:
        temp_list = eng_str_break(pro)
    else:
        temp_list = chi_str_break(pro)

    for x in temp_list:
        blank_slide_layout = prs.slide_layouts[11]
        slide = prs.slides.add_slide(blank_slide_layout)
        #add image for proverb_bg
        img_path = r'C:\Users\Wei Jie\PycharmProjects\H_World\venv\Proverb_proj\Sun_Svc_Proverbs_automation\Dependencies\Proverb_bg.png'
        pic = slide.shapes.add_picture(img_path, Inches(-0.47), Inches(4.64), Inches(14.2), Inches(2.87)) #path, left, top, width, height
        #add image for presider_bg
        img2_path = r'C:\Users\Wei Jie\PycharmProjects\H_World\venv\Proverb_proj\Sun_Svc_Proverbs_automation\Dependencies\Congre_bg.png'
        pic2 = slide.shapes.add_picture(img2_path, Inches(-0.45), Inches(1.21), Inches(8.15), Inches(4.11)) #path, left, top, width, height
        #fixed size for all 'presider' text boxes (standardized)
        txBox = slide.shapes.add_textbox(Inches(0.81), Inches(4.19), Inches(4.16), Inches(1.13)) #left, top, width, height
        tf = txBox.text_frame
        p1 = tf.paragraphs[0]
        ran = p1.add_run()
        if is_presider:
            ran.text = "Presider 司会"
        else:
            ran.text = "Congregation 会众"
        #font attributes
        font = ran.font
        font.name = 'Arial'
        font.size = Pt(40)
        font.bold = True
        font.color.rgb = RGBColor(255, 255, 0)
        #fixed size for all proverb text boxes (standardized)

        txBox2 = slide.shapes.add_textbox(Inches(0.63), Inches(5.08), Inches(12.05), Inches(2.14)) #left, top, width, height
        text_frame = txBox2.text_frame
        p = text_frame.paragraphs[0]
        run = p.add_run()
        run.text = x
        text_frame.word_wrap = True
        #font attributes
        font2 = run.font
        if is_english:
            font2.name = 'Arial'
        else:
            font2.name = 'Microsoft Yahei'
        font2.size = Pt(40)
        font2.color.rgb = RGBColor(255, 255, 255)
        font2.bold = True
    empty_list = []

#eng together proverb
eng_toget = eng_list[12].replace('Together', '', 1)
eng_toget = eng_toget.replace('-', '', 1)
eng_toget = eng_toget.replace('(Amen, hallelujah!)', '', 1)
eng_toget = eng_toget.replace('(Amen, Hallelujah!)', '', 1)
eng_toget = eng_toget.lstrip()
eng_toget = eng_toget.rstrip()
eng_toget = eng_str_break(eng_toget)
#chi together proverb
chi_toget = chi_list[12].replace('全體', '', 1)
chi_toget = chi_toget.replace('全体', '', 1)
chi_toget = chi_toget.replace('－', '', 1)
chi_toget = chi_toget.lstrip()
chi_toget = chi_toget.rstrip()
chi_toget = chi_str_break(chi_toget)


for item in eng_toget:
    blank_slide_layout = prs.slide_layouts[11]
    slide = prs.slides.add_slide(blank_slide_layout)
    # add image for proverb_bg
    img_path = r'C:\Users\Wei Jie\PycharmProjects\H_World\venv\Proverb_proj\Sun_Svc_Proverbs_automation\Dependencies\Proverb_bg.png'
    pic = slide.shapes.add_picture(img_path, Inches(-0.47), Inches(4.64), Inches(14.2),
                                   Inches(2.87))  # path, left, top, width, height
    # add image for presider_bg
    img2_path = r'C:\Users\Wei Jie\PycharmProjects\H_World\venv\Proverb_proj\Sun_Svc_Proverbs_automation\Dependencies\Congre_bg.png'
    pic2 = slide.shapes.add_picture(img2_path, Inches(-0.45), Inches(1.21), Inches(8.15),
                                    Inches(4.11))  # path, left, top, width, height
    # fixed size for all 'presider' text boxes (standardized)
    txBox = slide.shapes.add_textbox(Inches(0.81), Inches(4.19), Inches(4.16), Inches(1.13))  # left, top, width, height
    tf = txBox.text_frame
    p1 = tf.paragraphs[0]
    ran = p1.add_run()
    ran.text = "Together 全体"
    # font attributes
    font = ran.font
    font.name = 'Arial'
    font.size = Pt(40)
    font.bold = True
    font.color.rgb = RGBColor(255, 255, 0)
    # fixed size for all proverb text boxes (standardized)

    txBox2 = slide.shapes.add_textbox(Inches(0.63), Inches(5.08), Inches(12.05), Inches(2.14))  # left, top, width, height
    text_frame = txBox2.text_frame
    p = text_frame.paragraphs[0]
    run = p.add_run()
    run.text = item
    text_frame.word_wrap = True
    # font attributes
    font2 = run.font
    font2.name = 'Arial'
    font2.size = Pt(40)
    font2.color.rgb = RGBColor(255, 255, 255)
    font2.bold = True

for item in chi_toget:
    blank_slide_layout = prs.slide_layouts[11]
    slide = prs.slides.add_slide(blank_slide_layout)
    # add image for proverb_bg
    img_path = r'C:\Users\Wei Jie\PycharmProjects\H_World\venv\Proverb_proj\Sun_Svc_Proverbs_automation\Dependencies\Proverb_bg.png'
    pic = slide.shapes.add_picture(img_path, Inches(-0.47), Inches(4.64), Inches(14.2),
                                   Inches(2.87))  # path, left, top, width, height
    # add image for presider_bg
    img2_path = r'C:\Users\Wei Jie\PycharmProjects\H_World\venv\Proverb_proj\Sun_Svc_Proverbs_automation\Dependencies\Congre_bg.png'
    pic2 = slide.shapes.add_picture(img2_path, Inches(-0.45), Inches(1.21), Inches(8.15),
                                    Inches(4.11))  # path, left, top, width, height
    # fixed size for all 'presider' text boxes (standardized)
    txBox = slide.shapes.add_textbox(Inches(0.81), Inches(4.19), Inches(4.16), Inches(1.13))  # left, top, width, height
    tf = txBox.text_frame
    p1 = tf.paragraphs[0]
    ran = p1.add_run()
    ran.text = "Together 全体"
    # font attributes
    font = ran.font
    font.name = 'Arial'
    font.size = Pt(40)
    font.bold = True
    font.color.rgb = RGBColor(255, 255, 0)
    # fixed size for all proverb text boxes (standardized)

    txBox2 = slide.shapes.add_textbox(Inches(0.63), Inches(5.08), Inches(12.05), Inches(2.14))  # left, top, width, height
    text_frame = txBox2.text_frame
    p = text_frame.paragraphs[0]
    run = p.add_run()
    run.text = item
    text_frame.word_wrap = True
    # font attributes
    font2 = run.font
    font2.name = 'Microsoft Yahei'
    font2.size = Pt(40)
    font2.color.rgb = RGBColor(255, 255, 255)
    font2.bold = True

prs.save('Sun Svc P2 (Prov Rec)_import.pptx')
os.startfile('Sun Svc P2 (Prov Rec)_import.pptx')



